from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn
import requests
import io
import PyPDF2
from pptx import Presentation
import ollama

# ==========================================
# ⚙️ CONFIGURATION
# ==========================================
APPWRITE_PROJECT_ID = "69cda978001070a4493b"
APPWRITE_BUCKET_ID = "69cdaa6f002c333f4450"

app = FastAPI(title="AI Classroom Neural Engine")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class ChatRequest(BaseModel):
    fileId: str
    message: str

class QuizRequest(BaseModel):
    fileId: str

# New model for the Teacher's Grading request
class GradeRequest(BaseModel):
    assignment_desc: str

# ==========================================
# 🔍 HELPER: EXTRACT TEXT
# ==========================================
def extract_text_from_appwrite(file_id: str):
    url = f"https://cloud.appwrite.io/v1/storage/buckets/{APPWRITE_BUCKET_ID}/files/{file_id}/download?project={APPWRITE_PROJECT_ID}"
    response = requests.get(url)
    if response.status_code != 200:
        raise Exception(f"Appwrite Error: {response.status_code}")

    content = response.content
    text = ""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text = "\n".join(page.extract_text() for page in pdf_reader.pages if page.extract_text())
    except: pass

    if not text.strip():
        try:
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        except: pass
    return text

# ==========================================
# 🚀 ENDPOINTS
# ==========================================

@app.get("/")
def read_root():
    return {"status": "Online"}

@app.post("/ai/chat")
async def chat(req: ChatRequest):
    doc_text = extract_text_from_appwrite(req.fileId)
    res = ollama.chat(model='llama3', messages=[
        {'role': 'user', 'content': f"Context: {doc_text[:10000]}\n\nQuestion: {req.message}"}
    ])
    return {"reply": res['message']['content']}

@app.post("/ai/quiz")
async def quiz(req: QuizRequest):
    doc_text = extract_text_from_appwrite(req.fileId)
    res = ollama.chat(model='llama3', messages=[
        {'role': 'user', 'content': f"Generate a 3-question quiz from this text: {doc_text[:10000]}"}
    ])
    return {"quiz": res['message']['content']}

# 🔥 NEW: TEACHER AUTO-GRADER ENDPOINT
@app.post("/ai/grade-submission/{file_id}")
async def grade_submission(file_id: str, req: GradeRequest):
    print(f"⚖️ AI is grading student file: {file_id}")
    try:
        student_work = extract_text_from_appwrite(file_id)
        
        prompt = f"""
        You are a strict professor. Grade the student's submission based on this criteria:
        Criteria: {req.assignment_desc}
        
        Student Submission:
        {student_work[:12000]}
        
        Provide the output in this EXACT format:
        Grade: [X]/10
        Feedback: [Your short feedback here]
        """
        
        response = ollama.chat(model='llama3', messages=[
            {'role': 'user', 'content': prompt}
        ])
        
        return {"result": response['message']['content']}
    except Exception as e:
        return {"result": f"Grade: 0/10\nFeedback: Error processing file - {str(e)}"}

if __name__ == "__main__":
    uvicorn.run("main:app", host="127.0.0.1", port=8001, reload=True)